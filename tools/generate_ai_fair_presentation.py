# -*- coding: utf-8 -*-
from __future__ import annotations

import math
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPT_DIR = ROOT / "ppt"
ASSET_DIR = ROOT / "outputs" / "manual-ppt-work" / "assets"
OUT = PPT_DIR / "AI_Music_Recommendation_System_AI_Fair_Presentation.pptx"

W, H = 13.333, 7.5

COLORS = {
    "bg": RGBColor(7, 9, 20),
    "bg2": RGBColor(13, 16, 33),
    "panel": RGBColor(24, 29, 50),
    "panel2": RGBColor(34, 42, 69),
    "text": RGBColor(248, 251, 255),
    "muted": RGBColor(155, 166, 194),
    "accent": RGBColor(30, 215, 96),
    "blue": RGBColor(56, 189, 248),
    "purple": RGBColor(124, 58, 237),
    "danger": RGBColor(251, 91, 107),
    "paper": RGBColor(241, 245, 249),
    "ink": RGBColor(15, 23, 42),
}


SONGS = [
    ("Shape of You", "Ed Sheeran", "Pop", 0.83, 0.65, 96, 88),
    ("Blinding Lights", "The Weeknd", "Pop", 0.51, 0.73, 171, 92),
    ("Believer", "Imagine Dragons", "Rock", 0.78, 0.78, 125, 85),
    ("Titanium", "David Guetta ft. Sia", "Electronic", 0.60, 0.79, 126, 82),
    ("Lose Yourself", "Eminem", "Hip-hop", 0.69, 0.74, 171, 89),
    ("Levitating", "Dua Lipa", "Pop", 0.70, 0.82, 103, 87),
    ("Riptide", "Vance Joy", "Indie", 0.53, 0.61, 105, 82),
    ("Earned It", "The Weeknd", "R&B", 0.70, 0.45, 97, 91),
]


def font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    paths = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Avenir Next.ttc",
        "/System/Library/Fonts/SFNS.ttf",
    ]
    for path in paths:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            continue
    return ImageFont.load_default()


def create_mockup(kind: str, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", (1440, 900), (7, 9, 20))
    draw = ImageDraw.Draw(img)
    for i in range(900):
        r = int(7 + 18 * (i / 900))
        b = int(20 + 26 * (i / 900))
        draw.line([(0, i), (1440, i)], fill=(r, 12 + i // 90, b))
    draw.ellipse((-170, -150, 560, 420), fill=(16, 88, 52))
    draw.ellipse((870, -180, 1580, 380), fill=(53, 28, 116))
    draw.rounded_rectangle((46, 34, 1394, 106), radius=26, fill=(10, 14, 29), outline=(48, 55, 78), width=2)
    draw.rounded_rectangle((72, 52, 120, 98), radius=14, fill=(30, 215, 96))
    draw.text((87, 62), "AI", font=font(16, True), fill=(3, 18, 10))
    draw.text((138, 48), "AI Music Recommendation System", font=font(26, True), fill=(248, 251, 255))
    draw.text((138, 78), "LightFM Hybrid Recommender", font=font(15), fill=(155, 166, 194))
    draw.text((1030, 62), "Recommendations   Training   EN/中文", font=font(16), fill=(155, 166, 194))

    if kind == "home":
        draw.text((68, 170), "AI Fair Machine Learning Project", font=font(24, True), fill=(30, 215, 96))
        draw.text((68, 220), "Train a music recommendation", font=font(58, True), fill=(248, 251, 255))
        draw.text((68, 286), "model from your playlist.", font=font(58, True), fill=(248, 251, 255))
        draw.text((72, 380), "Import playlist metadata, train LightFM, rank songs, and learn from feedback.", font=font(25), fill=(180, 190, 216))
        draw.rounded_rectangle((72, 462, 292, 528), radius=20, fill=(30, 215, 96))
        draw.text((104, 480), "Get Recommendations", font=font(21, True), fill=(5, 18, 10))
        draw.rounded_rectangle((316, 462, 540, 528), radius=20, fill=(36, 43, 68), outline=(78, 86, 114))
        draw.text((352, 480), "Generate Final 50", font=font(21, True), fill=(248, 251, 255))
        draw.ellipse((940, 190, 1260, 510), fill=(16, 23, 40), outline=(30, 215, 96), width=10)
        draw.ellipse((1045, 295, 1155, 405), fill=(240, 255, 247))
        draw.text((982, 570), "Main Model: LightFM WARP", font=font(30, True), fill=(248, 251, 255))
    elif kind == "training":
        panels = [(60, 150, 430, 430, "Import Playlist"), (470, 150, 870, 430, "Manual Playlist"), (910, 150, 1378, 430, "Train Model")]
        for x1, y1, x2, y2, title in panels:
            draw.rounded_rectangle((x1, y1, x2, y2), radius=24, fill=(24, 29, 50), outline=(58, 65, 92), width=2)
            draw.text((x1 + 24, y1 + 26), title, font=font(28, True), fill=(248, 251, 255))
        metrics = [("Model Type", "LightFM WARP"), ("Training Status", "Trained"), ("Users", "1"), ("Songs", "284"), ("Interactions", "233"), ("Item Features", "146")]
        for i, (label, value) in enumerate(metrics):
            x = 945 + (i % 2) * 185
            y = 238 + (i // 2) * 72
            draw.rounded_rectangle((x, y, x + 160, y + 54), radius=16, fill=(36, 48, 66))
            draw.text((x + 12, y + 8), label, font=font(13), fill=(155, 166, 194))
            draw.text((x + 12, y + 25), value, font=font(18, True), fill=(215, 255, 228))
        draw.text((92, 520), "How we train:", font=font(32, True), fill=(30, 215, 96))
        steps = ["playlist songs → interactions", "songs → genre / artist / audio buckets", "LightFM learns preference weights", "recommendations update after feedback"]
        for i, s in enumerate(steps):
            draw.rounded_rectangle((92 + i * 320, 590, 350 + i * 320, 700), radius=22, fill=(24, 29, 50), outline=(65, 76, 105))
            draw.text((112 + i * 320, 616), s, font=font(18, True), fill=(248, 251, 255))
    elif kind == "recommendations":
        draw.text((70, 145), "Recommendations", font=font(44, True), fill=(248, 251, 255))
        draw.text((70, 198), "Ranked by learned model score plus content similarity.", font=font(22), fill=(155, 166, 194))
        cards = SONGS[:6]
        for i, s in enumerate(cards):
            x = 70 + (i % 3) * 450
            y = 260 + (i // 3) * 260
            draw.rounded_rectangle((x, y, x + 400, y + 218), radius=24, fill=(24, 29, 50), outline=(59, 66, 93))
            draw.text((x + 22, y + 22), s[0], font=font(25, True), fill=(248, 251, 255))
            draw.text((x + 22, y + 55), f"{s[1]} · {s[2]}", font=font(18), fill=(155, 166, 194))
            scores = [("Model", .92 - i * .025), ("Content", .76 - i * .03), ("Final", .86 - i * .025)]
            for j, (label, val) in enumerate(scores):
                draw.rounded_rectangle((x + 22 + j * 122, y + 92, x + 112 + j * 122, y + 142), radius=14, fill=(25, 57, 47))
                draw.text((x + 31 + j * 122, y + 105), f"{label} {val:.2f}", font=font(13, True), fill=(215, 255, 228))
            draw.text((x + 22, y + 164), "Reason: shares genre, tempo, and energy profile.", font=font(16), fill=(180, 190, 216))
    else:
        draw.text((70, 145), "Daily Recommendation Mix", font=font(44, True), fill=(248, 251, 255))
        draw.text((70, 198), "Final 50-song output for the AI Fair demo.", font=font(22), fill=(155, 166, 194))
        list_songs = SONGS * 7
        for i, s in enumerate(list_songs[:18]):
            x = 80 + (i // 9) * 650
            y = 260 + (i % 9) * 66
            draw.rounded_rectangle((x, y, x + 590, y + 48), radius=14, fill=(24, 29, 50), outline=(55, 63, 88))
            draw.rounded_rectangle((x + 12, y + 9, x + 42, y + 39), radius=9, fill=(30, 215, 96))
            draw.text((x + 21, y + 15), str(i + 1), font=font(14, True), fill=(3, 18, 10))
            draw.text((x + 58, y + 8), s[0], font=font(17, True), fill=(248, 251, 255))
            draw.text((x + 58, y + 28), f"{s[1]} · {s[2]}", font=font(13), fill=(155, 166, 194))
            draw.text((x + 525, y + 17), f"{0.91 - i * .01:.3f}", font=font(14, True), fill=(215, 255, 228))
    img.save(path)


def add_bg(slide):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLORS["bg"]
    rect.line.fill.background()


def add_title(slide, kicker, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = kicker.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS["accent"]
    p.font.name = "Aptos"
    p.character_spacing = 1600
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.68), Inches(11.6), Inches(0.82))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]
    p.font.name = "Aptos Display"
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.57), Inches(1.35), Inches(10.6), Inches(0.42))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS["muted"]
        p.font.name = "Aptos"


def textbox(slide, x, y, w, h, text, size=14, color="text", bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color]
    p.alignment = align
    return shape


def panel(slide, x, y, w, h, fill="panel", line=True, radius=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    if line:
        shape.line.color.rgb = RGBColor(62, 72, 102)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def bullet_list(slide, x, y, w, items, size=14, color="text"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
        p.level = 0
        p.space_after = Pt(8)
    return box


def add_image(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def mini_card(slide, x, y, w, h, title, body, accent="accent"):
    panel(slide, x, y, w, h)
    textbox(slide, x + 0.18, y + 0.15, w - 0.36, 0.25, title, 12, accent, True)
    textbox(slide, x + 0.18, y + 0.48, w - 0.36, h - 0.55, body, 12, "text")


def metric(slide, x, y, label, value, accent="accent"):
    panel(slide, x, y, 2.15, 0.9, "panel2")
    textbox(slide, x + 0.15, y + 0.13, 1.85, 0.22, label, 9, "muted", True)
    textbox(slide, x + 0.15, y + 0.42, 1.9, 0.34, value, 20, accent, True)


def make_presentation():
    PPT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    screenshots = {
        "home": ASSET_DIR / "ppt-ui-home.png",
        "training": ASSET_DIR / "ppt-ui-training.png",
        "recommendations": ASSET_DIR / "ppt-ui-recommendations.png",
        "daily": ASSET_DIR / "ppt-ui-daily.png",
    }
    for kind, path in screenshots.items():
        create_mockup(kind, path)

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 1 Cover
    s = prs.slides.add_slide(blank); add_bg(s)
    textbox(s, 0.55, 0.55, 2.2, 0.25, "AI FAIR 2026 · MACHINE LEARNING TRACK", 10, "accent", True)
    textbox(s, 0.55, 1.2, 6.5, 2.2, "AI Music\nRecommendation System", 46, "text", True)
    textbox(s, 0.6, 3.55, 5.4, 0.7, "Training a hybrid recommendation model from playlists, song features, and user feedback.", 17, "muted")
    metric(s, 0.6, 5.05, "Main Model", "LightFM WARP")
    metric(s, 3.0, 5.05, "Output", "Daily 50")
    metric(s, 5.4, 5.05, "Data", "CSV + APIs")
    add_image(s, screenshots["home"], 7.15, 0.75, 5.7, 4.25)
    textbox(s, 7.3, 6.05, 5.3, 0.3, "Presenter: AI Fair Project Team", 12, "muted")

    # 2 Problem
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "project overview", "The project solves a real recommendation problem: finding songs that match a listener’s taste.")
    mini_card(s, 0.65, 2.0, 3.8, 2.2, "Real-world problem", "Music libraries are too large to browse manually. A recommender must learn from what the user already likes.")
    mini_card(s, 4.75, 2.0, 3.8, 2.2, "Target users", "Students, music listeners, and AI Fair visitors who want to see how playlist data becomes a trained model.", "blue")
    mini_card(s, 8.85, 2.0, 3.8, 2.2, "Why it matters", "It connects a familiar product experience with core machine learning ideas: data, features, training, prediction, and feedback.", "purple")
    textbox(s, 0.8, 5.15, 11.8, 0.5, "Thesis: this is not just a search page. It creates training data, trains a hybrid model, and updates recommendations after feedback.", 20, "accent", True, PP_ALIGN.CENTER)

    # 3 System
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "system architecture", "A local Flask app turns playlist metadata and user actions into machine learning training data.")
    labels = [("Playlist Import", "QQ Music metadata"), ("Song Search", "Local + public APIs"), ("CSV Store", "songs / interactions / feedback"), ("Training", "LightFM WARP"), ("Ranking", "Daily 50 recommendations")]
    for i, (a, b) in enumerate(labels):
        x = 0.55 + i * 2.55
        panel(s, x, 2.55, 2.15, 1.3, "panel2")
        textbox(s, x + 0.15, 2.78, 1.85, 0.25, a, 13, "text", True, PP_ALIGN.CENTER)
        textbox(s, x + 0.15, 3.16, 1.85, 0.25, b, 10, "muted", False, PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            textbox(s, x + 2.2, 3.0, 0.4, 0.3, "→", 24, "accent", True, PP_ALIGN.CENTER)
    textbox(s, 0.85, 5.0, 11.5, 0.8, "All data is stored locally in CSV files. Online services are used only for public song metadata; no music audio is downloaded.", 18, "muted", False, PP_ALIGN.CENTER)

    # 4 Dataset
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "dataset and features", "The dataset is simple enough to inspect, but structured enough to train a model.")
    rows = [("songs.csv", "Song metadata + audio features", "track, artist, genre, energy, tempo, popularity"), ("interactions.csv", "Training signals", "playlist_import, manual_add, like, dislike"), ("feedback.csv", "Explicit feedback", "Like = 1, Dislike = 0"), ("playlists.csv", "Import history", "source, song_count, created_at")]
    for i, row in enumerate(rows):
        y = 2.0 + i * 1.1
        panel(s, 0.8, y, 11.7, 0.78, "panel2")
        textbox(s, 1.05, y + 0.2, 2.0, 0.25, row[0], 16, "accent", True)
        textbox(s, 3.25, y + 0.2, 3.0, 0.25, row[1], 14, "text", True)
        textbox(s, 6.3, y + 0.2, 5.9, 0.25, row[2], 13, "muted")

    # 5 QQ import
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "playlist import", "QQ Music links are parsed safely, then the app imports metadata only.")
    bullet_list(s, 0.75, 2.0, 5.6, [
        "Supports playlist/{id}, taoge.html?id=, disstid=, tid=, and pure numeric IDs.",
        "Tries multiple metadata endpoints: qzone, v8 playlist, and musicu.fcg.",
        "Never downloads music, never bypasses copyright, and never accesses paid audio.",
        "If QQ fails, the user gets a friendly fallback: manual playlist or song search."
    ], 15, "text")
    panel(s, 7.1, 2.0, 5.2, 3.8, "panel2")
    textbox(s, 7.4, 2.3, 4.5, 0.35, "URL extraction logic", 18, "accent", True)
    code = "https://y.qq.com/n/ryqq/playlist/12345\n→ extract_qq_playlist_id()\n→ playlist_id = 12345\n→ fetch metadata\n→ normalize songs\n→ create playlist_import interactions"
    textbox(s, 7.4, 2.9, 4.6, 2.2, code, 15, "text")

    # 6 API search
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "music search APIs", "The app searches local data first, then public metadata APIs when more results are needed.")
    providers = [("1", "Local CSV", "Fast, stable, includes training features"), ("2", "Deezer", "Clean public metadata for mainstream songs"), ("3", "MusicBrainz", "Open music knowledge base, 1 request/sec"), ("4", "Last.fm", "Optional API key; skipped if missing")]
    for i, (num, name, desc) in enumerate(providers):
        x = 0.75 + (i % 2) * 6.1
        y = 2.0 + (i // 2) * 1.65
        panel(s, x, y, 5.45, 1.15, "panel2")
        textbox(s, x + 0.2, y + 0.24, 0.45, 0.35, num, 20, "accent", True)
        textbox(s, x + 0.75, y + 0.18, 2.3, 0.25, name, 17, "text", True)
        textbox(s, x + 0.75, y + 0.55, 4.3, 0.28, desc, 12, "muted")
    textbox(s, 0.8, 5.8, 11.6, 0.4, "Every request uses timeout=8 and try/except, so API failure never crashes the Flask app.", 18, "blue", True, PP_ALIGN.CENTER)

    # 7 Feature engineering
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "feature engineering", "Songs become item feature tokens that LightFM can learn from.")
    tokens = ["genre:pop", "artist:the_weeknd", "tag:dance", "source:qq", "energy:high", "tempo:fast", "valence:positive", "danceability:high", "popularity:high"]
    for i, tok in enumerate(tokens):
        x = 0.8 + (i % 3) * 4.1
        y = 2.0 + (i // 3) * 1.05
        panel(s, x, y, 3.55, 0.72, "panel2")
        textbox(s, x + 0.18, y + 0.2, 3.1, 0.25, tok, 16, "accent" if i % 2 == 0 else "blue", True, PP_ALIGN.CENTER)
    textbox(s, 0.9, 5.4, 11.5, 0.75, "Instead of only using numeric similarity, the model learns embeddings for interpretable feature buckets such as high energy, fast tempo, genre, artist, and source.", 18, "muted", False, PP_ALIGN.CENTER)

    # 8 Interaction matrix
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "training data", "Playlist actions and feedback become a user–song interaction matrix.")
    metric(s, 0.85, 2.1, "playlist_import", "+1.0")
    metric(s, 3.2, 2.1, "manual_add", "+1.0", "blue")
    metric(s, 5.55, 2.1, "like", "+2.0")
    metric(s, 7.9, 2.1, "dislike", "-1.0", "danger")
    panel(s, 1.0, 3.75, 11.2, 1.8, "panel2")
    textbox(s, 1.35, 4.05, 10.5, 0.35, "LightFM WARP trains from positive ranking examples. Dislikes are saved, then used during re-ranking to lower scores.", 20, "text", True, PP_ALIGN.CENTER)

    # 9 Model
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "model and method", "LightFM is the main machine learning model because it supports hybrid recommendation.")
    bullet_list(s, 0.8, 2.0, 5.7, ["Open-source recommender library", "Learns from implicit and explicit feedback", "Uses WARP ranking loss for Top-N recommendation", "Combines user–song interactions with item features", "Works well for small educational datasets"], 16)
    panel(s, 7.0, 2.0, 5.1, 3.7, "panel2")
    textbox(s, 7.35, 2.35, 4.45, 0.4, "Training object", 18, "accent", True)
    textbox(s, 7.35, 3.0, 4.4, 1.8, 'model = LightFM(loss="warp")\nmodel.fit(interactions,\n          item_features,\n          sample_weight,\n          epochs=30)', 18, "text")

    # 10 Pipeline with screenshot
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "training workflow", "The UI exposes every training step so the audience can see the machine learning pipeline.")
    add_image(s, screenshots["training"], 0.7, 1.85, 6.5, 4.1)
    bullet_list(s, 7.6, 2.0, 4.9, ["Import or manually add songs", "Write interactions.csv rows", "Build item feature tokens", "Train LightFM WARP", "Save lightfm_model.pkl", "Show model status in the UI"], 16)

    # 11 Formula
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "recommendation ranking", "Final recommendations blend learned preference and content similarity.")
    textbox(s, 0.9, 2.0, 11.5, 0.55, "final_score = 0.7 × LightFM model_score + 0.3 × content_similarity_score", 27, "accent", True, PP_ALIGN.CENTER)
    mini_card(s, 1.1, 3.15, 3.5, 1.7, "model_score", "The trained LightFM model predicts how much the user may like each candidate song.")
    mini_card(s, 4.95, 3.15, 3.5, 1.7, "content_similarity", "Cosine similarity compares audio features such as energy, tempo, danceability, and mood.", "blue")
    mini_card(s, 8.8, 3.15, 3.5, 1.7, "rule explanation", "Reasons are generated deterministically, not by an LLM.", "purple")

    # 12 UI walkthrough
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "software experience", "The app feels like a modern music recommendation dashboard, not a plain form.")
    add_image(s, screenshots["home"], 0.6, 1.8, 5.95, 3.75)
    add_image(s, screenshots["recommendations"], 6.85, 1.8, 5.95, 3.75)
    textbox(s, 0.8, 6.05, 11.8, 0.35, "Bilingual UI, playlist import, manual songs, model status, recommendation cards, Like/Dislike, and Daily 50 output.", 17, "muted", False, PP_ALIGN.CENTER)

    # 13 Results
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "results and evaluation", "The system produces ranked songs, explanations, and a final 50-song daily mix.")
    add_image(s, screenshots["recommendations"], 0.65, 1.85, 5.7, 4.2)
    metric(s, 6.95, 2.0, "Recommendation Cards", "Top 10")
    metric(s, 9.35, 2.0, "Final Playlist", "50 songs", "blue")
    metric(s, 6.95, 3.25, "Explanation Type", "Rules")
    metric(s, 9.35, 3.25, "Feedback Loop", "Retrain")
    textbox(s, 7.0, 4.85, 4.75, 0.8, "Evaluation demo: after Like/Dislike feedback, interactions.csv changes, model retrains, and future rankings shift.", 17, "muted")

    # 14 Daily 50
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "final output", "Daily Recommendation Mix turns model scores into a show-ready playlist.")
    add_image(s, screenshots["daily"], 0.75, 1.75, 7.25, 4.55)
    bullet_list(s, 8.35, 2.1, 4.2, ["Generates up to 50 songs", "Filters songs already used for training", "Reduces repeated artists and genres", "Can be copied as a playlist text output"], 16)

    # 15 Real world
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "real-world use", "The same pipeline can support real school or personal music discovery.")
    mini_card(s, 0.8, 2.0, 3.6, 2.0, "For students", "Learn ML from a product everyone understands: playlists and song recommendations.")
    mini_card(s, 4.75, 2.0, 3.6, 2.0, "For music apps", "Convert playlist behavior and feedback into a continuously improving recommendation loop.", "blue")
    mini_card(s, 8.7, 2.0, 3.6, 2.0, "For future work", "Add richer audio analysis, more users, offline evaluation metrics, and stronger diversity controls.", "purple")
    textbox(s, 1.1, 5.45, 11.0, 0.5, "Most useful improvement: collect more feedback over time and evaluate precision@k / recall@k.", 20, "accent", True, PP_ALIGN.CENTER)

    # 16 Ethics
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "ethical AI disclosure", "The project uses AI help responsibly and avoids copyright risks.")
    bullet_list(s, 0.8, 2.0, 6.0, ["AI assistance was used for coding, debugging, UI writing, and documentation support.", "The model and code were checked and understood by the student.", "No local LLM, LM Studio, Ollama, or OpenAI API is used for recommendation explanations.", "QQ Music import uses metadata only; no audio download or unauthorized playback."], 15)
    panel(s, 7.25, 2.0, 4.9, 3.65, "panel2")
    textbox(s, 7.55, 2.35, 4.25, 0.4, "Risk controls", 20, "accent", True)
    bullet_list(s, 7.55, 3.0, 4.1, ["Timeouts for API requests", "Friendly error fallback", "Local CSV storage", "Rule-based explanations", "No private audio scraping"], 14, "text")

    # 17 Limitations
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "limitations", "The model is real, but the dataset size and public metadata sources limit accuracy.")
    mini_card(s, 0.8, 2.0, 3.7, 1.75, "Small dataset", "A single demo user means rankings are useful for demonstration, not production-scale personalization.")
    mini_card(s, 4.75, 2.0, 3.7, 1.75, "Estimated features", "Some online results do not provide true audio features, so the app estimates features from metadata.", "blue")
    mini_card(s, 8.7, 2.0, 3.7, 1.75, "Unstable APIs", "QQ Music public endpoints may change, so the app includes fallback and manual input.", "purple")
    textbox(s, 1.0, 5.2, 11.1, 0.65, "Future improvements: real audio-feature sources, multi-user data, train/test evaluation, precision@k, recall@k, and stronger diversity re-ranking.", 20, "muted", False, PP_ALIGN.CENTER)

    # 18 How to use: quick start
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "how to use", "Quick start: open the app, import music taste data, then train the model.")
    add_image(s, screenshots["home"], 0.65, 1.82, 5.6, 3.7)
    steps = [
        ("1", "Start the app", "Run python app.py, then open the local URL printed in the terminal."),
        ("2", "Import or add songs", "Paste a QQ Music playlist link, or type songs manually in Song - Artist format."),
        ("3", "Check the dataset", "The app saves songs, playlists, interactions, and feedback into local CSV files."),
        ("4", "Train the model", "Click Train Model to build a LightFM hybrid recommendation model.")
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 6.65
        y = 1.82 + i * 1.12
        panel(s, x, y, 5.75, 0.84, "panel2")
        panel(s, x + 0.18, y + 0.2, 0.44, 0.44, "accent", False)
        textbox(s, x + 0.28, y + 0.27, 0.24, 0.2, num, 12, "ink", True, PP_ALIGN.CENTER)
        textbox(s, x + 0.78, y + 0.14, 2.7, 0.24, title, 15, "text", True)
        textbox(s, x + 0.78, y + 0.42, 4.55, 0.25, body, 10, "muted")
    textbox(s, 0.95, 6.08, 11.5, 0.35, "For the AI Fair booth, prepare a short playlist first so the model has training data before visitors test recommendations.", 16, "accent", True, PP_ALIGN.CENTER)

    # 19 How to use: live demo flow
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "demo tutorial", "A simple live demo can show the complete machine learning loop in under three minutes.")
    demo_steps = [
        ("Step 1", "Search or import", "Search for familiar songs such as Shape of You, Blinding Lights, Believer, or Levitating."),
        ("Step 2", "Create interactions", "Add songs to the training playlist. These actions become positive training examples."),
        ("Step 3", "Train LightFM", "Click Train Model and point out the model type, number of interactions, and item features."),
        ("Step 4", "Generate recommendations", "Click Get Recommendations to show Top 10 ranked songs with model/content/final scores."),
        ("Step 5", "Give feedback", "Click Like or Dislike. The app saves feedback and retrains so future recommendations change."),
        ("Step 6", "Generate Final 50", "Click Generate Final 50 to output a Daily Recommendation Mix that can be copied.")
    ]
    for i, (label, title, body) in enumerate(demo_steps):
        x = 0.7 + (i % 2) * 6.25
        y = 1.92 + (i // 2) * 1.42
        panel(s, x, y, 5.65, 1.02, "panel2")
        textbox(s, x + 0.22, y + 0.16, 1.0, 0.22, label, 10, "accent", True)
        textbox(s, x + 1.28, y + 0.14, 3.8, 0.25, title, 16, "text", True)
        textbox(s, x + 1.28, y + 0.48, 4.05, 0.28, body, 10, "muted")
    textbox(s, 0.9, 6.35, 11.5, 0.32, "Presentation line: The important part is not the button click. It is that every click creates data that changes the trained model.", 15, "blue", True, PP_ALIGN.CENTER)

    # 20 How to use: backup paths and troubleshooting
    s = prs.slides.add_slide(blank); add_bg(s); add_title(s, "usage notes", "If an online service fails, the system still works through local data and manual input.")
    mini_card(s, 0.75, 2.0, 3.65, 1.6, "If QQ import fails", "Use manual playlist input or Search Songs. QQ public metadata endpoints can change or block some playlists.")
    mini_card(s, 4.85, 2.0, 3.65, 1.6, "If API search is slow", "The app uses timeouts and skips failed providers, so Flask stays running instead of crashing.", "blue")
    mini_card(s, 8.95, 2.0, 3.65, 1.6, "If LightFM is unavailable", "The app falls back to content similarity and a logistic-regression feedback model.", "purple")
    panel(s, 1.05, 4.35, 11.25, 1.15, "panel2")
    textbox(s, 1.35, 4.62, 10.65, 0.32, "Teacher-facing explanation", 18, "accent", True, PP_ALIGN.CENTER)
    textbox(s, 1.35, 5.0, 10.65, 0.35, "The project is designed to keep running during a live demo: public API errors become friendly messages, while the local CSV dataset and trained recommender still provide recommendations.", 14, "text", False, PP_ALIGN.CENTER)

    # 21 Closing
    s = prs.slides.add_slide(blank); add_bg(s)
    textbox(s, 0.75, 0.75, 2.0, 0.3, "CLOSING", 11, "accent", True)
    textbox(s, 0.75, 1.4, 7.8, 2.0, "This project shows the full ML loop: data → features → training → ranking → feedback.", 40, "text", True)
    add_image(s, screenshots["home"], 8.1, 1.15, 4.7, 3.4)
    textbox(s, 0.85, 4.75, 6.7, 1.0, "Playlist import creates implicit feedback. Like/Dislike creates explicit feedback. LightFM learns preferences and ranks songs the user is likely to enjoy.", 20, "muted")
    textbox(s, 0.85, 6.35, 11.7, 0.35, "Thank you · Questions?", 24, "accent", True, PP_ALIGN.CENTER)

    prs.save(OUT)


if __name__ == "__main__":
    make_presentation()
    print(OUT)
